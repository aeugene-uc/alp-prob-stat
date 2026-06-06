# -*- coding: utf-8 -*-
import json, warnings, io, contextlib
warnings.filterwarnings('ignore')
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
plt.show = lambda *a, **k: None
import numpy as np
import pandas as pd

# ---------- 1. Re-execute notebook logic to capture real tables ----------
nb = json.load(open('AI.ipynb', encoding='utf-8'))
ns = {}
def clean(src):
    out = []
    for ln in src.split('\n'):
        if 'plotly' in ln or ln.strip().startswith('pio.'):
            continue
        out.append(ln)
    return '\n'.join(out)

buf = io.StringIO()
with contextlib.redirect_stdout(buf):
    for c in nb['cells']:
        if c['cell_type'] != 'code':
            continue
        src = clean(''.join(c['source']))
        try:
            exec(src, ns)
        except Exception as e:
            print('CELLERR', repr(e))
        plt.close('all')

errs = [l for l in buf.getvalue().split('\n') if l.startswith('CELLERR')]
print('exec done. cell errors:', errs if errs else 'none')

# Pull captured objects
df          = ns['df']
raw         = ns['raw']
pearson_eff = ns['pearson_eff']
spearman_eff= ns['spearman_eff']
perf        = ns['perf']
report      = ns['report']
se_cmp      = ns['se_cmp']
abla        = ns['abla']
abla_year   = ns['abla_year']
banding_mod = ns['banding']          # reassigned in cell55 = model comparison
gain        = ns['gain']
bp_p        = ns['bp_p']; jb_p = ns['jb_p']
skew        = ns['skew']; kurt = ns['kurt']
exkurt      = ns['exkurt']; R2_aux = ns['R2_aux']
m_yr        = ns['m_yr']
r2t_full    = ns['r2t_full']; r2t_nc = ns['r2t_nc']; r2t_year = ns['r2t_year']

# Descriptive table (numerics incl year + target) reconstructed
desc_src = pd.DataFrame({
    'engine_size': raw['ENGINE SIZE'].values,
    'cylinders':   raw['CYLINDERS'].values,
    'year':        raw['YEAR'].values,
    'n_gears':     df['n_gears'].values,
    'comb_l_100km':raw['COMB (L/100 km)'].values,
})
desc = desc_src.describe().T[['count', 'mean', 'std', 'min', '50%', 'max']].round(2)

# Pearson vs Spearman comparison (rebuild; banding var was overwritten)
ps = pd.DataFrame({'Pearson_r': pearson_eff['r'], 'Spearman_rho': spearman_eff['rho']})
ps['selisih'] = (ps['Spearman_rho'] - ps['Pearson_r'])
ps = ps.sort_values('selisih', key=abs, ascending=False).round(4)

# r2 summary
import sklearn.metrics as skm
r2_tr = report  # placeholder not used

print('captured tables OK')

# ---------- 2. Build the PPTX ----------
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

INK   = RGBColor(0x21, 0x21, 0x21)
GRAY  = RGBColor(0x5F, 0x5F, 0x5F)
LIGHT = RGBColor(0xEC, 0xEC, 0xEC)
HEAD  = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE  = RGBColor(0xBB, 0xBB, 0xBB)
FONT  = 'Calibri'

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def _set_run(r, size, color=INK, bold=False, italic=False):
    r.font.name = FONT; r.font.size = Pt(size)
    r.font.color.rgb = color; r.font.bold = bold; r.font.italic = italic

def title(s, text, sub=None):
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.42), Inches(12.1), Inches(0.95))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    _set_run(r, 30, INK, bold=True)
    if sub:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = sub
        _set_run(r2, 14, GRAY)
    # rule line
    ln = s.shapes.add_shape(1, Inches(0.62), Inches(1.42), Inches(12.1), Pt(2))
    ln.fill.solid(); ln.fill.fore_color.rgb = RULE; ln.line.fill.background()
    ln.shadow.inherit = False
    return tb

def bullets(s, items, left, top, width, height, size=16, gap=6):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        lvl = 0
        txt = it
        if it.startswith('>'):
            lvl = 1; txt = it[1:].strip()
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.space_before = Pt(0)
        r = p.add_run()
        r.text = ('     –  ' if lvl else '•  ') + txt
        _set_run(r, size - (1 if lvl else 0), GRAY if lvl else INK)
    return tb

def image(s, path, left, top, max_w, max_h, center=True):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = max_w; h = w / ar
    if h > max_h:
        h = max_h; w = h * ar
    l = left + (max_w - w) / 2 if center else left
    t = top + (max_h - h) / 2 if center else top
    s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

NOSTYLE = '{5940675A-B579-460E-94D1-54222C63F5DA}'  # Table Grid, no banding
def table(s, df, left, top, width, fontsize=12, index_name='', col_fmt=None,
          colw=None, rowh=0.32, header_h=0.36):
    df = df.copy()
    nrows = len(df) + 1
    ncols = len(df.columns) + 1
    gx = s.shapes.add_table(nrows, ncols, Inches(left), Inches(top),
                            Inches(width), Inches(header_h + rowh * len(df)))
    tbl = gx.table
    # neutral style
    tbl.first_row = False; tbl.horz_banding = False
    tbl._tbl.tblPr.set('firstRow', '0'); tbl._tbl.tblPr.set('bandRow', '0')
    for el in tbl._tbl.tblPr.findall(qn('a:tableStyleId')):
        tbl._tbl.tblPr.remove(el)
    sid = tbl._tbl.tblPr.makeelement(qn('a:tableStyleId'), {}); sid.text = NOSTYLE
    tbl._tbl.tblPr.append(sid)
    # column widths
    if colw is None:
        colw = [1.4] + [1.0] * (ncols - 1)
    tot = sum(colw)
    for j, cw in enumerate(colw):
        tbl.columns[j].width = Emu(int(Inches(width) * cw / tot))
    # header
    hdr = [index_name] + list(df.columns)
    for j, txt in enumerate(hdr):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = HEAD
        cell.margin_left = Inches(0.05); cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
        rr = p.add_run(); rr.text = str(txt); _set_run(rr, fontsize, WHITE, bold=True)
    # body
    for i in range(len(df)):
        for j in range(ncols):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT
            cell.margin_left = Inches(0.05); cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if j == 0:
                val = str(df.index[i])
            else:
                v = df.iloc[i, j - 1]
                col = df.columns[j - 1]
                if col_fmt and col in col_fmt:
                    val = col_fmt[col](v)
                elif isinstance(v, (int, np.integer)):
                    val = str(int(v))
                elif isinstance(v, (float, np.floating)):
                    val = f'{v:.4f}' if abs(v) < 1 else f'{v:.3f}'
                else:
                    val = str(v)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            rr = p.add_run(); rr.text = val; _set_run(rr, fontsize, INK)
    return gx

def footer(s, n):
    tb = s.shapes.add_textbox(Inches(11.9), Inches(7.02), Inches(1.3), Inches(0.35))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f'ALP Stats  |  {n}'; _set_run(r, 9, RULE)

IMG = '_slides_img'
pidx = [0]
def newnum():
    pidx[0] += 1; return pidx[0]

# ===== Slide 1 — Title =====
s = slide()
bg = s.background.fill; bg.solid(); bg.fore_color.rgb = WHITE
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run(); r.text = 'ALP Stats'
_set_run(r, 54, INK, bold=True)
p2 = tf.add_paragraph(); r = p2.add_run()
r.text = 'Pemodelan Konsumsi Bahan Bakar Kendaraan'
_set_run(r, 24, INK); p2.space_before = Pt(10)
p3 = tf.add_paragraph(); r = p3.add_run()
r.text = 'Regresi Linear Berganda + One-Hot Encoding'
_set_run(r, 16, GRAY)
ln = s.shapes.add_shape(1, Inches(0.95), Inches(4.45), Inches(5.5), Pt(2))
ln.fill.solid(); ln.fill.fore_color.rgb = RULE; ln.line.fill.background(); ln.shadow.inherit = False
names = [('Rei Putra Soemanto', '0706012410060'),
         ('Amadeus Eugene Dirgantara', '0706012410063'),
         ('Jason Tio', '0706012410006')]
tb2 = s.shapes.add_textbox(Inches(0.95), Inches(4.7), Inches(11), Inches(1.8))
tf2 = tb2.text_frame
for i, (nm, nim) in enumerate(names):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    r = p.add_run(); r.text = nm + '   '; _set_run(r, 16, INK, bold=True)
    r2 = p.add_run(); r2.text = nim; _set_run(r2, 14, GRAY)
    p.space_after = Pt(4)

# ===== Slide 2 — Dataset & Tujuan =====
s = slide(); title(s, 'Dataset & Tujuan')
bullets(s, [
    'Sumber: Fuel Consumption 2000-2022; 22.556 baris kendaraan',
    'Target: comb_l_100km (konsumsi bahan bakar gabungan, L/100 km)',
    'Prediktor: spesifikasi mesin, kelas kendaraan, bahan bakar, transmisi',
    'Tujuan: memodelkan & menjelaskan konsumsi dari spesifikasi kendaraan',
    'Metode: regresi linear berganda dengan one-hot encoding (OLS)',
    'Pendekatan inferensi: keputusan berbasis ukuran efek & confidence interval',
    '>p-value tidak informatif pada n besar (~22 ribu): hampir selalu signifikan',
], 0.7, 1.8, 12.0, 5.0, size=18, gap=10)
footer(s, newnum())

# ===== Slide 3 — Pra-pemrosesan =====
s = slide(); title(s, 'Pra-pemrosesan Data')
bullets(s, [
    'Nilai hilang: tidak ada di kolom mana pun -> tanpa imputasi',
    '1 baris duplikat persis dihapus -> tersisa 22.555 observasi unik',
    'Label vehicle_class dinormalkan (kapitalisasi & pemisah - / :)',
    '>32 label berbeda menyusut menjadi 17 kelas yang konsisten',
    'Kode fuel & transmisi ditulis sebagai "definisi (kode)"',
    '>mis. Regular gasoline (X), Automatic (A)',
    'transmission dipecah menjadi transmission_type + n_gears',
    '>CVT (AV) tidak punya gigi tetap -> n_gears = 0 (konvensi)',
], 0.7, 1.8, 12.0, 5.0, size=18, gap=8)
footer(s, newnum())

# ===== Slide 4 — Variabel & Peran =====
s = slide(); title(s, 'Variabel & Peran')
roles = pd.DataFrame(
    {'Tipe': ['numerik', 'numerik', 'numerik', 'numerik',
              'kategorikal', 'kategorikal', 'kategorikal', 'target'],
     'Peran': ['prediktor', 'prediktor', 'dibuang (korelasi ~0)', 'prediktor',
               'prediktor (one-hot)', 'prediktor (one-hot)', 'prediktor (one-hot)', 'respons'],
     'Kategori/Catatan': ['kontinu, 1-8 L', '4 / 6 / 8', 'model year 2000-2022', '0 (CVT) - 10',
                          '5 jenis', '17 kelas (dirapikan)', '5 tipe', 'L/100 km']},
    index=['engine_size', 'cylinders', 'year', 'n_gears',
           'fuel', 'vehicle_class', 'transmission_type', 'comb_l_100km'])
table(s, roles, 0.7, 1.75, 12.0, fontsize=13, index_name='Variabel',
      colw=[1.5, 1.0, 2.0, 2.6], rowh=0.46, header_h=0.4)
footer(s, newnum())

# ===== Slide 5 — Statistik Deskriptif =====
s = slide(); title(s, 'Statistik Deskriptif')
table(s, desc, 0.7, 1.8, 9.2, fontsize=14, index_name='',
      colw=[1.6, 0.9, 0.9, 0.9, 0.9, 0.9, 0.9], rowh=0.5, header_h=0.42)
bullets(s, [
    'Mesin 1-8 L; silinder 4/6/8',
    'year merata 2000-2022',
    'n_gears 0 (CVT) - 10',
    'Dominan Regular & Premium gasoline',
], 10.2, 1.9, 2.9, 4.5, size=13, gap=8)
footer(s, newnum())

# ===== Slide 6 — Pair Plot =====
s = slide(); title(s, 'Pair Plot Variabel Numerik')
image(s, f'{IMG}/plot1_cell26.png', 0.6, 1.55, 8.0, 5.5)
bullets(s, [
    'engine_size, cylinders, comb_l_100km menanjak jelas bersama',
    'cylinders & n_gears diskret -> pita-pita vertikal',
    'year hampir datar terhadap konsumsi',
    'Hubungan dengan target mendekati garis lurus',
    '>mendukung korelasi Pearson & model linear',
], 8.8, 1.9, 4.3, 5.0, size=15, gap=10)
footer(s, newnum())

# ===== Slide 7 — Distribusi per Bahan Bakar =====
s = slide(); title(s, 'Distribusi Konsumsi per Bahan Bakar')
image(s, f'{IMG}/plot2_cell29.png', 0.6, 1.6, 8.2, 5.2)
bullets(s, [
    'Median bergeser sistematis antar kategori',
    'Diesel paling hemat',
    'Ethanol E85 & Natural Gas paling boros',
    '>densitas energi lebih rendah',
    'Bahan bakar layak masuk model',
], 9.0, 1.9, 4.1, 5.0, size=15, gap=10)
footer(s, newnum())

# ===== Slide 8 — Distribusi per Kelas Kendaraan =====
s = slide(); title(s, 'Distribusi Konsumsi per Kelas Kendaraan')
image(s, f'{IMG}/plot3_cell29.png', 0.6, 1.6, 8.4, 5.4)
bullets(s, [
    'Gradien jelas dari mobil kecil ke besar',
    'Sedan & kompak paling hemat',
    'Van & pickup paling boros',
    'Kelas kendaraan membawa informasi kuat',
], 9.1, 2.0, 4.0, 5.0, size=15, gap=10)
footer(s, newnum())

# ===== Slide 9 — Distribusi per Transmisi =====
s = slide(); title(s, 'Distribusi Konsumsi per Transmisi')
image(s, f'{IMG}/plot4_cell29.png', 0.6, 1.6, 8.2, 5.2)
bullets(s, [
    'Perbedaan antar tipe lebih landai namun sistematis',
    'CVT cenderung paling hemat',
    'Lima tipe transmisi tetap informatif',
    'Layak di-one-hot encode',
], 9.0, 2.0, 4.1, 5.0, size=15, gap=10)
footer(s, newnum())

# ===== Slide 10 — Korelasi Pearson =====
s = slide(); title(s, 'Korelasi Pearson', 'Hubungan linear; keputusan via CI 95% (Fisher z) terhadap ambang |r| > 0.1')
image(s, f'{IMG}/plot5_cell32.png', 0.5, 1.7, 5.4, 4.7)
pe = pearson_eff[['r', 'CI_bawah', 'CI_atas', 'Putusan']].round(3)
table(s, pe, 6.1, 1.75, 7.0, fontsize=11, index_name='Pasangan',
      colw=[2.3, 0.9, 1.1, 1.1, 1.5], rowh=0.46, header_h=0.4)
bullets(s, [
    'engine_size r=0.81 & cylinders r=0.77 -> kuat',
    'engine_size~cylinders r=0.91 -> kolinier',
    'year r=-0.07 -> dapat diabaikan',
], 6.1, 5.3, 7.0, 1.6, size=12, gap=5)
footer(s, newnum())

# ===== Slide 11 — Korelasi Spearman =====
s = slide(); title(s, 'Korelasi Spearman', 'Hubungan monoton via peringkat; robust outlier; CI 95% Bonett-Wright (2000)')
image(s, f'{IMG}/plot6_cell37.png', 0.5, 1.7, 5.4, 4.7)
sp = spearman_eff[['rho', 'CI_bawah', 'CI_atas', 'Putusan']].round(3)
table(s, sp, 6.1, 1.75, 7.0, fontsize=11, index_name='Pasangan',
      colw=[2.3, 0.9, 1.1, 1.1, 1.5], rowh=0.46, header_h=0.4)
bullets(s, [
    'engine_size 0.84 & cylinders 0.82 -> kuat',
    'engine_size~cylinders 0.94 -> kolinier',
    'pasangan ~year di batas -> tidak meyakinkan',
], 6.1, 5.3, 7.0, 1.6, size=12, gap=5)
footer(s, newnum())

# ===== Slide 12 — Pearson vs Spearman =====
s = slide(); title(s, 'Perbandingan Pearson vs Spearman', 'Penyaring linearitas informal (bukan uji formal)')
psd = ps.rename(columns={'Pearson_r': 'Pearson (r)', 'Spearman_rho': 'Spearman (rho)',
                         'selisih': 'selisih (S - P)'})
table(s, psd, 0.7, 1.85, 8.4, fontsize=12, index_name='Pasangan',
      colw=[2.2, 1.1, 1.1, 1.2], rowh=0.46, header_h=0.4)
bullets(s, [
    'Pearson = linear; Spearman = monoton',
    'Selisih kecil & positif pada pasangan target (+0.04, +0.05)',
    '>kelengkungan monoton sangat ringan',
    'Pearson sedikit underestimate, tetap sah',
    'Bentuk non-monoton tak tampak di keduanya',
    '>linearitas dicek via residual & model linear vs polinomial',
], 9.3, 1.95, 3.8, 5.0, size=13, gap=8)
footer(s, newnum())

# ===== Slide 13 — Seleksi Fitur & Encoding =====
s = slide(); title(s, 'Seleksi Fitur & One-Hot Encoding')
bullets(s, [
    'engine_size & cylinders dipertahankan',
    '>kolinier (r=0.91) tetapi tidak redundan; ditafsirkan bersama',
    'year dibuang: korelasi ~0 dan kontribusi prediktif ~0',
    'n_gears dipertahankan (transmisi berkaitan dengan konsumsi)',
    'fuel, vehicle_class, transmission_type di-one-hot (k-1 kolom)',
    'Kategori referensi: Diesel (D), COMPACT, Automatic (A)',
    '>menghindari dummy variable trap',
    'engine_size_sq disiapkan untuk uji linear vs polinomial',
], 0.7, 1.8, 12.2, 5.0, size=18, gap=8)
footer(s, newnum())

# ===== Slide 14 — Pemilihan Model & Performa =====
s = slide(); title(s, 'Pemilihan Model & Performa', 'Split 80/20 (random_state=1); model dipilih dari R2 test')
bm = banding_mod.copy()
bm.index = ['Linear', 'Polinomial (eng_size^2)']
bm = bm[['R2_train', 'R2_test', 'n_param']].round(4)
table(s, bm, 0.7, 1.85, 7.2, fontsize=12, index_name='Model',
      colw=[2.2, 1.1, 1.1, 1.0], rowh=0.5, header_h=0.4)
pf = perf.round(4)
table(s, pf, 0.7, 3.9, 7.2, fontsize=12, index_name='Metrik',
      colw=[1.6, 1.1, 1.1, 1.4], rowh=0.5, header_h=0.4)
bullets(s, [
    f'Kenaikan R2 test polinomial: {gain:+.4f} -> diabaikan',
    'Dipilih model LINEAR (parsimoni)',
    'R2 train & test berdekatan -> tidak overfit',
    'RMSE & MAE dalam satuan asli (L/100 km)',
    'F-test: prediktor berpengaruh bersama (p~0)',
], 8.4, 1.95, 4.7, 5.0, size=14, gap=10)
footer(s, newnum())

# ===== Slide 15 — Interpretasi Koefisien =====
s = slide(); title(s, 'Interpretasi Koefisien', 'Koefisien & CI 95% berbasis standard error robust HC3')
keys = ['engine_size', 'cylinders', 'n_gears',
        'fuel_Ethanol E85 (E)', 'fuel_Natural gas (N)',
        'fuel_Premium gasoline (Z)', 'fuel_Regular gasoline (X)',
        'transmission_type_Continuously variable (AV)',
        'vehicle_class_VAN PASSENGER', 'vehicle_class_VAN CARGO']
labels = ['engine_size', 'cylinders', 'n_gears',
          'fuel: Ethanol E85', 'fuel: Natural gas', 'fuel: Premium', 'fuel: Regular',
          'transmisi: CVT', 'kelas: VAN PASSENGER', 'kelas: VAN CARGO']
ct = report.loc[keys, ['coef', 'CI_low_95', 'CI_high_95', 'p_value']].copy()
ct.index = labels
ct = ct.round(3)
table(s, ct, 0.7, 1.85, 8.6, fontsize=12, index_name='Prediktor (vs referensi)',
      colw=[2.4, 1.0, 1.1, 1.1, 1.0], rowh=0.44, header_h=0.4)
bullets(s, [
    'Mesin lebih besar & silinder lebih banyak -> boros',
    'Tiap gigi tambahan sedikit lebih hemat',
    'Diesel paling hemat; E85 paling boros',
    'CVT paling hemat antar transmisi',
    'Van & pickup paling boros antar kelas',
], 9.5, 1.95, 3.6, 5.0, size=13, gap=9)
footer(s, newnum())

# ===== Slide 16 — Validasi Asumsi =====
s = slide(); title(s, 'Validasi Asumsi', 'Pada n besar p-value selalu menolak; dinilai dari MAGNITUDO')
asm = pd.DataFrame(
    {'Nilai': [f'{bp_p:.1e}', f'{jb_p:.1e}', f'{skew:+.2f}', f'{exkurt:+.2f}', f'{R2_aux:.3f}'],
     'Ambang / makna': ['tolak homoskedastik', 'tolak normal',
                        '|.|<=0.5 -> hampir simetris', '|.|<=1 -> ekor berat (leptokurtik)',
                        'proporsi var. residual^2 dijelaskan']},
    index=['Breusch-Pagan p', 'Jarque-Bera p', 'Skewness',
           'Excess kurtosis', 'BP R2_aux'])
table(s, asm, 0.7, 1.85, 8.6, fontsize=12, index_name='Uji / Ukuran',
      colw=[1.9, 1.2, 3.0], rowh=0.5, header_h=0.4)
bullets(s, [
    'Heteroskedastisitas nyata (sedang)',
    '>ditangani SE robust HC3',
    'Ekor berat = penyimpangan utama',
    'Inferensi koefisien aman (CLT)',
    'Interval prediksi kurang terkalibrasi di ekor',
], 9.5, 1.95, 3.6, 5.0, size=13, gap=9)
footer(s, newnum())

# ===== Slide 17 — Diagnostik Visual =====
s = slide(); title(s, 'Diagnostik Visual Residual')
image(s, f'{IMG}/plot7_cell71.png', 0.5, 1.7, 8.6, 4.8)
bullets(s, [
    'Q-Q: tengah lurus, ekor menyimpang',
    '>leptokurtik (ekor berat)',
    'Residual vs fitted: pola corong',
    '>heteroskedastisitas (sesuai Breusch-Pagan)',
    'Tak ada lengkungan sistematis',
    '>spesifikasi linear memadai',
    'Ditangani dengan SE robust HC3',
], 9.2, 1.9, 3.9, 5.2, size=14, gap=8)
footer(s, newnum())

# ===== Slide 18 — Refinement: Ablasi & HC3 =====
s = slide(); title(s, 'Refinement: Ablasi & Standard Error Robust (HC3)')
ab = pd.DataFrame(
    {'R2_test': [r2t_full, r2t_nc, r2t_year],
     'Delta vs final': [0.0, r2t_nc - r2t_full, r2t_year - r2t_full]},
    index=['Model final', 'Tanpa cylinders', 'Tambah year']).round(4)
table(s, ab, 0.7, 1.85, 6.0, fontsize=12, index_name='Skenario',
      colw=[2.0, 1.2, 1.4], rowh=0.5, header_h=0.4)
sec = se_cmp.loc[['engine_size', 'cylinders', 'fuel_Premium gasoline (Z)',
                  'vehicle_class_VAN PASSENGER']].copy()
sec.index = ['engine_size', 'cylinders', 'fuel: Premium', 'kelas: VAN PASS.']
sec = sec.round(4)
table(s, sec, 0.7, 4.35, 6.6, fontsize=11, index_name='Koefisien (sampel)',
      colw=[2.0, 1.0, 1.0, 1.2], rowh=0.46, header_h=0.4)
bullets(s, [
    'Buang cylinders: R2 test turun ~0.02',
    '>kolinier bukan redundan -> dipertahankan',
    f'Tambah year: R2 test {r2t_year - r2t_full:+.4f}, p={m_yr.pvalues["year"]:.2f}',
    '>keputusan membuang year terkonfirmasi',
    'HC3: SE engine_size & cylinders naik ~40%',
    '>estimasi titik tetap, ketidakpastian dikoreksi',
], 7.6, 1.95, 5.5, 5.0, size=14, gap=9)
footer(s, newnum())

# ===== Slide 19 — Kesimpulan =====
s = slide(); title(s, 'Kesimpulan')
bullets(s, [
    'Model menjelaskan ~84% variansi konsumsi; generalisasi baik (tanpa overfitting)',
    'Penggerak utama: jenis bahan bakar, kelas kendaraan, ukuran mesin',
    '>n_gears & tipe transmisi berpengaruh lebih kecil (CVT paling hemat)',
    'engine_size^2 tidak menambah R2 test -> model linear (parsimoni)',
    'Heteroskedastisitas dikoreksi dengan SE robust HC3',
    'Non-normalitas: inferensi koefisien tetap valid (CLT)',
    '>interval prediksi kurang terkalibrasi di ekor',
    'Keterbatasan: kemungkinan korelasi dalam-grup; model asosiatif, bukan kausal',
], 0.7, 1.8, 12.3, 5.0, size=18, gap=9)
footer(s, newnum())

prs.save('ALP_Stats.pptx')
print('saved ALP_Stats.pptx with', len(prs.slides._sldIdLst), 'slides')
