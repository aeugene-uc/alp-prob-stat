# -*- coding: utf-8 -*-
"""Deck hitam-putih bergaya paper ALP_Paper_Revisi.pdf.
Times New Roman, tabel booktabs (garis horizontal saja), plot grayscale.
Maks 20 slide. Tanpa ';' dan tanpa '-'/'->' sebagai penyambung kalimat.
"""
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

BLACK = RGBColor(0x00, 0x00, 0x00)
DGRAY = RGBColor(0x3A, 0x3A, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = 'Times New Roman'

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
NOGRID = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'  # No Style, No Grid

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = WHITE
    return s

def _set_run(r, size, color=BLACK, bold=False, italic=False):
    r.font.name = FONT; r.font.size = Pt(size)
    r.font.color.rgb = color; r.font.bold = bold; r.font.italic = italic

def title(s, text, sub=None):
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.40), Inches(12.1), Inches(0.95))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    _set_run(r, 27, BLACK, bold=True)
    if sub:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = sub
        _set_run(r2, 13, DGRAY, italic=True)
    ln = s.shapes.add_shape(1, Inches(0.62), Inches(1.36 if not sub else 1.58), Inches(12.1), Pt(1.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLACK; ln.line.fill.background()
    ln.shadow.inherit = False
    return tb

def bullets(s, items, left, top, width, height, size=16, gap=8):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        lvl = 0; txt = it
        if it.startswith('>'):
            lvl = 1; txt = it[1:].strip()
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.space_before = Pt(0)
        r = p.add_run()
        r.text = ('       '  + txt) if lvl else ('•  ' + txt)
        _set_run(r, size - (1 if lvl else 0), DGRAY if lvl else BLACK, italic=bool(lvl))
    return tb

def image(s, path, left, top, max_w, max_h, center=True):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = max_w; h = w / ar
    if h > max_h:
        h = max_h; w = h * ar
    l = left + (max_w - w) / 2 if center else left
    t = top + (max_h - h) / 2 if center else top
    s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

def caption(s, text, left, top, width, size=11):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; _set_run(r, size, DGRAY, italic=True)

def _set_border(cell, edge, width_emu):
    tag = qn('a:ln' + edge)
    tcPr = cell._tc.get_or_add_tcPr()
    for e in tcPr.findall(tag):
        tcPr.remove(e)
    ln = tcPr.makeelement(tag, {'w': str(width_emu), 'cap': 'flat', 'cmpd': 'sng', 'algn': 'ctr'})
    sf = ln.makeelement(qn('a:solidFill'), {})
    clr = sf.makeelement(qn('a:srgbClr'), {'val': '000000'})
    sf.append(clr); ln.append(sf)
    tcPr.insert(0, ln)  # lines must precede fill in schema order

THICK = 19050   # 1.5 pt
THIN  = 9525    # 0.75 pt

def booktable(s, df, left, top, width, fontsize=12, index_name='', colw=None,
              rowh=0.34, header_h=0.4, aligns=None):
    nrows = len(df) + 1
    ncols = len(df.columns) + 1
    gx = s.shapes.add_table(nrows, ncols, Inches(left), Inches(top),
                            Inches(width), Inches(header_h + rowh * len(df)))
    tbl = gx.table
    tbl.first_row = False; tbl.horz_banding = False
    tbl._tbl.tblPr.set('firstRow', '0'); tbl._tbl.tblPr.set('bandRow', '0')
    for el in tbl._tbl.tblPr.findall(qn('a:tableStyleId')):
        tbl._tbl.tblPr.remove(el)
    sid = tbl._tbl.tblPr.makeelement(qn('a:tableStyleId'), {}); sid.text = NOGRID
    tbl._tbl.tblPr.append(sid)
    if colw is None:
        colw = [1.6] + [1.0] * (ncols - 1)
    tot = sum(colw)
    for j, cw in enumerate(colw):
        tbl.columns[j].width = Emu(int(Inches(width) * cw / tot))
    if aligns is None:
        aligns = ['l'] + ['r'] * (ncols - 1)
    almap = {'l': PP_ALIGN.LEFT, 'c': PP_ALIGN.CENTER, 'r': PP_ALIGN.RIGHT}

    def fill_cell(cell, txt, bold, align):
        cell.fill.background()
        cell.margin_left = Inches(0.07); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.015); cell.margin_bottom = Inches(0.015)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = txt; _set_run(r, fontsize, BLACK, bold=bold)

    hdr = [index_name] + list(df.columns)
    for j, txt in enumerate(hdr):
        fill_cell(tbl.cell(0, j), str(txt), True, almap[aligns[j]])
    for i in range(len(df)):
        for j in range(ncols):
            val = str(df.index[i]) if j == 0 else str(df.iloc[i, j - 1])
            fill_cell(tbl.cell(i + 1, j), val, False, almap[aligns[j]])
    # booktabs rules: top of header, bottom of header, bottom of last row
    for j in range(ncols):
        _set_border(tbl.cell(0, j), 'T', THICK)
        _set_border(tbl.cell(0, j), 'B', THIN)
        _set_border(tbl.cell(nrows - 1, j), 'B', THICK)
    return gx

def footer(s):
    # footer dihilangkan sesuai permintaan
    return

def captbl(s, text, left, top, width):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text; _set_run(r, 12, BLACK, bold=True)

FIG = 'fig_ai/'  # plot berwarna

# ============================ SLIDE 1 — TITLE ============================
s = slide()
tb = s.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11.3), Inches(2.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'Pemodelan Konsumsi Bahan Bakar Gabungan Kendaraan'
_set_run(r, 30, BLACK, bold=True)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'Menggunakan Regresi Linear Berganda dengan One-Hot Encoding'
_set_run(r, 30, BLACK, bold=True)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(8)
r = p.add_run(); r.text = 'Modeling Combined Vehicle Fuel Consumption Using Multiple Linear Regression with One-Hot Encoding'
_set_run(r, 14, DGRAY, italic=True)
ln = s.shapes.add_shape(1, Inches(4.4), Inches(4.05), Inches(4.5), Pt(1.2))
ln.fill.solid(); ln.fill.fore_color.rgb = BLACK; ln.line.fill.background(); ln.shadow.inherit = False
tb2 = s.shapes.add_textbox(Inches(1.0), Inches(4.25), Inches(11.3), Inches(2.2))
tf2 = tb2.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'Rei Putra Soemanto, Amadeus Eugene Dirgantara, Jason Tio'
_set_run(r, 17, BLACK, bold=True)
for line, sz, it in [
    ('NIM: 0706012410060, 0706012410063, 0706012410006', 13, False),
    ('Program Studi Informatika, Universitas Ciputra, Surabaya, Jawa Timur, Indonesia', 13, False),
    ('Tugas Mata Kuliah Probabilitas dan Statistika', 13, True)]:
    p = tf2.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(4)
    r = p.add_run(); r.text = line; _set_run(r, sz, DGRAY, italic=it)

# ============================ SLIDE 2 — Abstrak ============================
s = slide(); title(s, 'Abstrak')
bullets(s, [
    'Memodelkan konsumsi bahan bakar gabungan (comb_l_100km) dari spesifikasi kendaraan',
    'Dataset Fuel Consumption 2000 sampai 2022, 22.555 baris setelah pembersihan',
    'Metode regresi linear berganda dengan one-hot encoding, diestimasi OLS',
    'Tahun model dibuang, jumlah silinder dipertahankan karena ablasi membuktikan manfaatnya',
    'Suku kuadratik tidak menaikkan R-squared uji sehingga model linear dipilih',
    'Model menjelaskan sekitar 84 persen variansi dengan RMSE uji sekitar 1,18 L/100 km',
    'Heteroskedastisitas dan non-normalitas ditangani dengan standard error robust HC3',
], 0.7, 1.7, 12.0, 4.4, size=17, gap=10)
bullets(s, [
    'Kata kunci. regresi linear berganda, one-hot encoding, seleksi fitur, multikolinieritas, '
    'heteroskedastisitas, standard error robust HC3',
], 0.7, 6.35, 12.0, 0.8, size=13, gap=2)
footer(s)

# ============================ SLIDE 3 — Pendahuluan ============================
s = slide(); title(s, '1. Pendahuluan')
bullets(s, [
    'Konsumsi bahan bakar berkaitan langsung dengan biaya, emisi, dan keberlanjutan',
    'Memprediksinya dari spesifikasi kendaraan bermanfaat untuk armada dan konsumen',
    'Karakteristik kendaraan bercampur antara besaran numerik dan kategori',
    '>kategori diolah menjadi kolom indikator melalui one-hot encoding',
    'Tantangan utama dalam pemodelan ini',
    '>multikolinieritas antara ukuran mesin dan jumlah silinder',
    '>prediktor hampir tidak relevan seperti tahun model',
    '>bentuk hubungan yang mungkin melengkung',
    '>pemenuhan asumsi Ordinary Least Squares',
], 0.7, 1.7, 12.2, 5.3, size=17, gap=7)
footer(s)

# ============================ SLIDE 4 — Pertanyaan Penelitian ============================
s = slide(); title(s, 'Pertanyaan Penelitian')
bullets(s, [
    'Pertama, fitur kendaraan mana yang sebaiknya dipertahankan sebagai prediktor',
    'Kedua, bagaimana variabel kategori diolah tanpa menimbulkan kolinieritas sempurna',
    'Ketiga, apakah model linear sudah memadai atau perlu suku polinomial',
    'Keempat, apakah model akhir memenuhi asumsi OLS dan bagaimana menangani pelanggarannya',
], 0.7, 2.0, 12.0, 4.5, size=19, gap=20)
footer(s)

# ============================ SLIDE 5 — Tinjauan Pustaka ============================
s = slide(); title(s, '2. Tinjauan Pustaka')
bullets(s, [
    'Regresi linear berganda memodelkan respons sebagai jumlah kontribusi linear tiap prediktor',
    'One-hot encoding membuat k dikurang satu kolom indikator dengan satu kategori acuan',
    '>penetapan acuan menghindari dummy variable trap',
    'Pearson mengukur hubungan linear, Spearman mengukur hubungan monoton pada peringkat',
    'Multikolinieritas membesarkan standard error namun tidak selalu menurunkan daya prediksi',
    'Asumsi OLS diuji dengan Breusch-Pagan dan Jarque-Bera serta diagnostik visual',
    '>standard error robust HC3 dipakai saat ada heteroskedastisitas',
], 0.7, 1.7, 12.2, 5.3, size=17, gap=9)
footer(s)

# ============================ SLIDE 6 — Metode: Dataset & Variabel ============================
s = slide(); title(s, '3. Metode: Dataset dan Variabel')
roles = pd.DataFrame(
    {'Tipe': ['Numerik', 'Numerik', 'Numerik', 'Numerik', 'Kategori', 'Kategori', 'Kategori', 'Target'],
     'Peran': ['Prediktor', 'Prediktor', 'Dibuang (korelasi nol)', 'Prediktor',
               'Prediktor (one-hot)', 'Prediktor (one-hot)', 'Prediktor (one-hot)', 'Respons'],
     'Catatan': ['Kontinu, 0,8 sampai 8,4 L', 'Umumnya 4, 6, 8', 'Model 2000 sampai 2022',
                 'Nol untuk CVT sampai 10', '5 jenis', '17 kelas (dirapikan)', '5 tipe', 'L/100 km']},
    index=['Ukuran mesin', 'Jumlah silinder', 'Tahun model', 'Jumlah gigi',
           'Jenis bahan bakar', 'Kelas kendaraan', 'Tipe transmisi', 'comb_l_100km'])
booktable(s, roles, 0.7, 1.75, 12.0, fontsize=14, index_name='Variabel',
          colw=[1.7, 1.0, 1.9, 2.6], rowh=0.52, header_h=0.45,
          aligns=['l', 'l', 'l', 'l'])
footer(s)

# ============================ SLIDE 7 — Metode: Persiapan & Kerangka ============================
s = slide(); title(s, '3. Metode: Persiapan Data dan Kerangka Analisis')
bullets(s, [
    'Tidak ada nilai hilang, satu baris duplikat dihapus sehingga tersisa 22.555 baris',
    'Kelas kendaraan diseragamkan dari 32 menjadi 17 kelas',
    'Transmisi dipecah menjadi tipe transmisi dan jumlah gigi (CVT diberi nol gigi)',
    'Data dibagi 80 berbanding 20 menjadi 18.044 latih dan 4.511 uji',
], 0.7, 1.65, 12.2, 2.2, size=16, gap=7)
bullets(s, [
    'Tahap 1 eksplorasi data',
    'Tahap 2 seleksi fitur berbasis korelasi dan uji ablasi',
    'Tahap 3 pembentukan model (one-hot, split, linear vs polinomial)',
    'Tahap 4 validasi asumsi (Breusch-Pagan dan Jarque-Bera)',
    'Tahap 5 penyempurnaan inferensi dengan standard error robust HC3',
], 0.9, 3.7, 12.0, 3.2, size=17, gap=9)
footer(s)

# ============================ SLIDE 8 — Tabel 1 Deskriptif ============================
s = slide(); title(s, '4.1 Statistik Deskriptif')
captbl(s, 'Tabel 1. Statistik deskriptif variabel numerik (n = 22.555)', 0.7, 1.55, 12.0)
t1 = pd.DataFrame({
    'Mean': ['3,357', '5,854', '2011,6', '5,635', '11,034'],
    'Std': ['1,335', '1,820', '6,298', '1,957', '2,911'],
    'Min': ['0,8', '2', '2000', '0', '3,6'],
    'Q1': ['2,3', '4', '2006', '5', '9,1'],
    'Median': ['3,0', '6', '2012', '6', '10,6'],
    'Q3': ['4,2', '8', '2017', '6', '12,7'],
    'Max': ['8,4', '16', '2022', '10', '26,1'],
}, index=['Ukuran mesin (L)', 'Jumlah silinder', 'Tahun model', 'Jumlah gigi', 'comb_l_100km'])
booktable(s, t1, 0.7, 1.95, 12.0, fontsize=14, index_name='Variabel',
          colw=[2.0, 1.0, 1.0, 0.9, 0.9, 1.0, 0.9, 0.9], rowh=0.55, header_h=0.48)
bullets(s, [
    'Ukuran mesin median 3,0 L, silinder terpusat di 4, 6, dan 8',
    'Respons condong ke kanan dengan ekor melampaui 20 L/100 km',
], 0.7, 5.6, 12.0, 1.2, size=15, gap=8)
footer(s)

# ============================ SLIDE 9 — EDA Pair Plot ============================
s = slide(); title(s, '4.2 Eksplorasi Visual: Pair Plot')
image(s, FIG + 'pairplot.png', 0.5, 1.6, 7.4, 5.2)
caption(s, 'Gambar 1. Pair plot variabel numerik (sampel acak 3.000 baris)', 0.5, 6.75, 7.4, size=10)
bullets(s, [
    'Ukuran mesin, jumlah silinder, dan konsumsi menanjak jelas bersama',
    'Jumlah silinder dan jumlah gigi diskret sehingga membentuk pola vertikal',
    'Tahun model hampir datar terhadap konsumsi',
    'Hubungan dengan respons mendekati garis lurus',
    '>mendukung pemakaian model linear',
], 8.2, 1.9, 4.8, 5.0, size=15, gap=11)
footer(s)

# ============================ SLIDE 10 — EDA Boxplot ============================
s = slide(); title(s, '4.2 Eksplorasi Visual: Boxplot per Kategori')
image(s, FIG + 'box_vehicle_class.png', 0.45, 1.65, 6.5, 5.0, center=False)
caption(s, 'Gambar 3. Konsumsi menurut kelas kendaraan', 0.45, 6.7, 6.5, size=9)
image(s, FIG + 'box_fuel.png', 7.2, 1.7, 5.9, 2.0, center=False)
caption(s, 'Gambar 2. menurut jenis bahan bakar', 7.2, 3.5, 5.9, size=9)
image(s, FIG + 'box_transmission.png', 7.2, 4.0, 5.9, 2.0, center=False)
caption(s, 'Gambar 4. menurut tipe transmisi', 7.2, 5.8, 5.9, size=9)
bullets(s, [
    'Median konsumsi bergeser sistematis antar kategori',
    'Diesel paling hemat, Ethanol E85 dan Natural gas paling boros',
    'Kelas menanjak dari sedan kecil ke van dan pikap',
], 7.2, 6.05, 5.9, 1.3, size=12, gap=4)
footer(s)

# ============================ SLIDE 11 — Tabel 2 Korelasi ============================
s = slide(); title(s, '4.3 Korelasi dan Seleksi Fitur')
image(s, FIG + 'pearson.png', 0.5, 1.7, 5.0, 4.3)
caption(s, 'Gambar A1. Matriks korelasi Pearson', 0.5, 6.05, 5.0, size=10)
captbl(s, 'Tabel 2. Korelasi prediktor numerik', 5.8, 1.6, 7.3)
t2 = pd.DataFrame({
    'Pearson': ['0,807', '0,772', '-0,068', '0,913'],
    'Spearman': ['0,845', '0,818', '-0,073', '0,936'],
    'Putusan': ['Berarti', 'Berarti', 'Dapat diabaikan', 'Kolinier'],
}, index=['Ukuran mesin ~ respons', 'Jumlah silinder ~ respons',
          'Tahun model ~ respons', 'Ukuran mesin ~ silinder'])
booktable(s, t2, 5.8, 1.95, 7.3, fontsize=12.5, index_name='Pasangan',
          colw=[2.6, 1.0, 1.0, 1.6], rowh=0.5, header_h=0.45,
          aligns=['l', 'r', 'r', 'c'])
bullets(s, [
    'Putusan dari selang kepercayaan 95 persen terhadap ambang 0,1',
    'Pada n besar p-value hampir selalu signifikan sehingga tidak dipakai',
    'Selisih Spearman dan Pearson kecil sehingga model linear wajar',
], 5.8, 4.3, 7.3, 2.3, size=13, gap=9)
footer(s)

# ============================ SLIDE 12 — Tabel 3 Seleksi ============================
s = slide(); title(s, '4.3 Keputusan Seleksi Fitur')
captbl(s, 'Tabel 3. Keputusan seleksi fitur', 0.7, 1.55, 12.0)
t3 = pd.DataFrame({
    'Keputusan': ['Pertahankan', 'Pertahankan', 'Buang', 'Pertahankan',
                  'Pertahankan', 'Pertahankan', 'Pertahankan'],
    'Alasan ringkas': ['Korelasi paling kuat dengan respons',
                       'Kolinier namun terbukti menambah prediksi (ablasi)',
                       'Korelasi dengan respons dapat diabaikan',
                       'Transmisi berkaitan dengan konsumsi',
                       'Median konsumsi bergeser antar kategori',
                       'Median konsumsi bergeser antar kategori',
                       'Median konsumsi bergeser antar kategori']},
    index=['Ukuran mesin', 'Jumlah silinder', 'Tahun model', 'Jumlah gigi',
           'Jenis bahan bakar', 'Kelas kendaraan', 'Tipe transmisi'])
booktable(s, t3, 0.7, 1.95, 12.0, fontsize=14, index_name='Prediktor',
          colw=[1.7, 1.3, 4.0], rowh=0.56, header_h=0.45, aligns=['l', 'l', 'l'])
footer(s)

# ============================ SLIDE 13 — Tabel 4 Linear vs Polinomial ============================
s = slide(); title(s, '4.4 Pembentukan Model dan Pemilihan Bentuk')
captbl(s, 'Tabel 4. Perbandingan model linear dan polinomial', 0.7, 1.7, 11.0)
t4 = pd.DataFrame({
    'Jumlah parameter': ['28', '29'],
    'R-squared latih': ['0,8342', '0,8344'],
    'R-squared uji': ['0,8365', '0,8364'],
}, index=['Linear', 'Polinomial (ukuran mesin kuadratik)'])
booktable(s, t4, 0.7, 2.1, 11.0, fontsize=15, index_name='Model',
          colw=[3.0, 1.4, 1.4, 1.4], rowh=0.62, header_h=0.5, aligns=['l', 'r', 'r', 'r'])
bullets(s, [
    'Pada data uji kedua model praktis setara padahal polinomial satu parameter lebih',
    'Tambahan suku kuadratik tidak menaikkan kemampuan generalisasi',
    '>model linear yang lebih sederhana dipilih (parsimoni)',
], 0.7, 4.3, 12.0, 2.0, size=16, gap=10)
footer(s)

# ============================ SLIDE 14 — Tabel 5 Koefisien ============================
s = slide(); title(s, '4.5 Model Final dan Koefisien',
                    'Cuplikan. Standard error, p-value, dan CI 95 persen berbasis HC3. Tabel penuh ada di paper')
captbl(s, 'Tabel 5. Koefisien model linear final dengan standard error robust HC3', 0.7, 1.75, 12.4)
t5 = pd.DataFrame({
    'Koef': ['0,6438', '0,5740', '-0,0518', '5,4887', '3,7854', '2,0951', '1,4221',
             '-1,6222', '3,9195', '2,6157', '1,5188', '-0,1952'],
    'SE (HC3)': ['0,0264', '0,0185', '0,0075', '0,0758', '0,2542', '0,0574', '0,0569',
                 '0,0646', '0,1053', '0,0718', '0,0365', '0,0441'],
    'p-value': ['<0,001', '<0,001', '<0,001', '<0,001', '<0,001', '<0,001', '<0,001',
                '<0,001', '<0,001', '<0,001', '<0,001', '<0,001'],
    'CI bawah': ['0,592', '0,538', '-0,066', '5,340', '3,287', '1,983', '1,311',
                 '-1,749', '3,713', '2,475', '1,447', '-0,282'],
    'CI atas': ['0,696', '0,610', '-0,037', '5,637', '4,284', '2,208', '1,534',
                '-1,496', '4,126', '2,757', '1,590', '-0,109'],
}, index=['Ukuran mesin', 'Jumlah silinder', 'Jumlah gigi',
          'Bahan bakar: Ethanol E85 (E)', 'Bahan bakar: Natural gas (N)',
          'Bahan bakar: Premium gasoline (Z)', 'Bahan bakar: Regular gasoline (X)',
          'Transmisi: Continuously variable (AV)', 'Kelas: Van passenger',
          'Kelas: Van cargo', 'Kelas: SUV', 'Kelas: Minicompact'])
booktable(s, t5, 0.7, 2.15, 12.4, fontsize=11.5, index_name='Variabel (vs kategori acuan)',
          colw=[3.1, 0.95, 0.95, 0.95, 0.95, 0.95], rowh=0.355, header_h=0.4,
          aligns=['l', 'r', 'r', 'r', 'r', 'r'])
footer(s)

# ============================ SLIDE 15 — Interpretasi ============================
s = slide(); title(s, '4.5 Interpretasi Koefisien')
bullets(s, [
    'Ukuran mesin dan jumlah silinder positif sehingga mesin lebih besar lebih boros',
    '>keduanya dibaca bersama sebagai ukuran kapasitas mesin',
    'Jumlah gigi negatif sehingga lebih banyak rasio gigi cenderung lebih hemat',
    'Diesel paling hemat sebagai acuan, Ethanol E85 paling boros',
    'Transmisi variabel kontinu paling hemat antar tipe transmisi',
    'Kelas kendaraan menanjak dari mobil kecil sampai van dan pikap',
    '>kelas Full size dan Mid size tidak berbeda nyata dari acuan Compact',
], 0.7, 1.7, 12.2, 5.2, size=17, gap=10)
footer(s)

# ============================ SLIDE 16 — Tabel 6 Asumsi + diagnostik ============================
s = slide(); title(s, '4.6 Validasi Asumsi')
captbl(s, 'Tabel 6. Hasil uji asumsi dan magnitudo penyimpangan', 0.6, 1.55, 12.2)
t6 = pd.DataFrame({
    'Uji dan magnitudo': ['Breusch-Pagan, R-squared aux 0,083', 'Skewness 0,42', 'Excess kurtosis 2,14'],
    'Pembacaan': ['Heteroskedastik ringan sampai sedang', 'Kemiringan dapat diabaikan',
                  'Ekor lebih tebal dari normal'],
}, index=['Homoskedastisitas', 'Normalitas (kemiringan)', 'Normalitas (ekor)'])
booktable(s, t6, 0.6, 1.95, 7.2, fontsize=11.5, index_name='Aspek',
          colw=[1.7, 2.6, 2.4], rowh=0.5, header_h=0.45, aligns=['l', 'l', 'l'])
bullets(s, [
    'Pada n besar uji formal selalu menolak sehingga dibaca dari magnitudo',
    'Heteroskedastisitas ditangani standard error robust HC3',
    'Inferensi koefisien tetap aman berkat Central Limit Theorem',
], 0.6, 3.9, 7.2, 2.6, size=13, gap=10)
image(s, FIG + 'diagnostic.png', 8.0, 2.4, 5.1, 3.3)
caption(s, 'Gambar 5. Q-Q plot (kiri) dan residual vs fitted (kanan)', 8.0, 5.55, 5.1, size=10)
footer(s)

# ============================ SLIDE 17 — Tabel 7 HC3 ============================
s = slide(); title(s, '4.6 Penyempurnaan Inferensi HC3')
captbl(s, 'Tabel 7. Perbandingan standard error klasik dan HC3 (sampel koefisien)', 0.7, 1.7, 11.5)
t7 = pd.DataFrame({
    'SE klasik': ['0,0185', '0,0131', '0,0782', '0,0859'],
    'SE HC3': ['0,0264', '0,0185', '0,0574', '0,1053'],
    'Rasio HC3 / klasik': ['1,43', '1,41', '0,73', '1,23'],
}, index=['Ukuran mesin', 'Jumlah silinder',
          'Bahan bakar: Premium gasoline (Z)', 'Kelas: Van passenger'])
booktable(s, t7, 0.7, 2.1, 11.5, fontsize=14, index_name='Koefisien (sampel)',
          colw=[3.0, 1.2, 1.2, 1.6], rowh=0.56, header_h=0.48, aligns=['l', 'r', 'r', 'r'])
bullets(s, [
    'Estimasi titik koefisien tidak berubah, hanya ketidakpastiannya dikoreksi',
    'Standard error ukuran mesin dan jumlah silinder naik sekitar 40 persen',
    '>bukti standard error klasik kurang tepat saat ada heteroskedastisitas',
    'HC3 dipilih karena paling konservatif',
], 0.7, 4.5, 12.2, 2.2, size=15, gap=9)
footer(s)

# ============================ SLIDE 18 — Tabel 8 Ablasi & Generalisasi ============================
s = slide(); title(s, '4.7 Ablasi dan Evaluasi Generalisasi')
captbl(s, 'Tabel 8. Evaluasi generalisasi model final', 0.7, 1.6, 11.0)
t8 = pd.DataFrame({
    'Latih': ['0,8342', '1,1857', '0,8962'],
    'Uji': ['0,8365', '1,1757', '0,8941'],
    'Selisih (latih dikurang uji)': ['-0,0023', '0,0101', '0,0020'],
}, index=['R-squared', 'RMSE (L/100 km)', 'MAE (L/100 km)'])
booktable(s, t8, 0.7, 2.0, 11.0, fontsize=14, index_name='Metrik',
          colw=[2.2, 1.2, 1.2, 2.2], rowh=0.56, header_h=0.48, aligns=['l', 'r', 'r', 'r'])
bullets(s, [
    'Membuang jumlah silinder menurunkan R-squared uji sekitar 0,02',
    '>kolinier bukan berarti redundan sehingga dipertahankan',
    'Menambahkan kembali tahun model tidak mengubah R-squared uji (p sekitar 0,852)',
    'Selisih latih dan uji sangat kecil sehingga tidak ada overfitting',
], 0.7, 4.1, 12.2, 2.5, size=15, gap=9)
footer(s)

# ============================ SLIDE 19 — Kesimpulan ============================
s = slide(); title(s, '5. Kesimpulan')
bullets(s, [
    'Model menjelaskan sekitar 84 persen variansi konsumsi dan menggeneralisasi dengan baik',
    'Penggerak utama adalah jenis bahan bakar, kelas kendaraan, dan ukuran mesin',
    'Jumlah silinder dipertahankan karena ablasi membuktikan ia menambah daya prediksi',
    'Tahun model dibuang karena korelasi dan kontribusinya dapat diabaikan',
    'Suku kuadratik tidak diperlukan sehingga model linear yang dipilih',
    'Heteroskedastisitas dikoreksi dengan standard error robust HC3',
], 0.7, 1.7, 12.3, 5.2, size=17, gap=11)
footer(s)

# ============================ SLIDE 20 — Saran & Keterbatasan ============================
s = slide(); title(s, 'Keterbatasan, Saran, dan Sumber Data')
bullets(s, [
    'Breusch-Pagan dan Jarque-Bera menolak homoskedastisitas dan normalitas secara formal',
    '>penolakan wajar pada sampel besar dan tidak membatalkan estimasi titik koefisien',
    'Interval prediksi kurang terkalibrasi di bagian ekor akibat ekor tebal',
    'Independensi tidak diuji sehingga kemungkinan korelasi dalam grup masih terbuka',
    'Model bersifat asosiatif, bukan kausal',
    'Saran: menambah prediktor lain, memakai standard error berbasis cluster, kalibrasi interval prediksi',
    'Sumber data: Yilmaz, Fuel Consumption 2000 sampai 2022, Kaggle',
], 0.7, 1.7, 12.3, 5.2, size=16, gap=9)
footer(s)

OUT = 'ALP Statistics & Probability.pptx'
prs.save(OUT)
print('saved', OUT, 'with', len(prs.slides._sldIdLst), 'slides')
